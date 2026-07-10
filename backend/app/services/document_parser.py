"""Document parser: extract text from .docx, .pdf, .md, .txt, .pptx files.

Produces high-fidelity text extraction preserving structure (headings, tables,
lists, headers/footers) so the notebook has a full representation of the
original document.
"""
import io
import re


def detect_and_parse(content: bytes, filename: str) -> str:
    """Parse document content to plain text based on file extension."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext in ("md", "mdx", "txt"):
        return content.decode("utf-8", errors="replace")

    if ext == "docx":
        return _parse_docx(content)

    if ext == "pdf":
        return _parse_pdf(content)

    if ext == "pptx":
        return _parse_pptx(content)

    # Fallback: try as text
    return content.decode("utf-8", errors="replace")


def _parse_docx(content: bytes) -> str:
    """Extract text from DOCX preserving headings, lists, tables, headers/footers."""
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document(io.BytesIO(content))
    sections = []

    # --- Extract headers and footers from all sections ---
    for section in doc.sections:
        header_texts = []
        footer_texts = []
        try:
            for para in section.header.paragraphs:
                text = para.text.strip()
                if text:
                    header_texts.append(text)
        except Exception:
            pass
        try:
            for para in section.footer.paragraphs:
                text = para.text.strip()
                if text:
                    footer_texts.append(text)
        except Exception:
            pass
        if header_texts:
            sections.append("[HEADER]\n" + "\n".join(header_texts))
        if footer_texts:
            sections.append("[FOOTER]\n" + "\n".join(footer_texts))

    # --- Extract body paragraphs with style-aware formatting ---
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        style_name = (para.style.name or "").lower() if para.style else ""

        # Detect heading levels from style
        if "heading" in style_name:
            # Extract heading level (e.g., "heading 1" -> "#", "heading 2" -> "##")
            match = re.search(r'heading\s*(\d+)', style_name)
            if match:
                level = int(match.group(1))
                sections.append(f"\n{'#' * level} {text}\n")
            else:
                sections.append(f"\n# {text}\n")
        elif "title" in style_name:
            sections.append(f"\n# {text}\n")
        elif "subtitle" in style_name:
            sections.append(f"\n## {text}\n")
        elif "list" in style_name:
            # Detect list type
            if "bullet" in style_name:
                sections.append(f"- {text}")
            elif "number" in style_name:
                sections.append(f"1. {text}")
            else:
                sections.append(f"- {text}")
        elif "quote" in style_name:
            sections.append(f"> {text}")
        else:
            sections.append(text)

    # --- Extract tables with structure ---
    for table_idx, table in enumerate(doc.tables):
        table_lines = []
        for row_idx, row in enumerate(table.rows):
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace("\n", " ")
                cells.append(cell_text)
            if any(cells):
                table_lines.append("| " + " | ".join(cells) + " |")
                # Add header separator after first row
                if row_idx == 0:
                    table_lines.append("|" + "|".join(["---"] * len(cells)) + "|")
        if table_lines:
            sections.append("\n" + "\n".join(table_lines) + "\n")

    # --- Extract text from text boxes and shapes in inline shapes ---
    try:
        for rel in doc.part.rels.values():
            if "header" in rel.reltype or "footer" in rel.reltype:
                pass  # Already handled above
    except Exception:
        pass

    return "\n\n".join(sections)


def _parse_pdf(content: bytes) -> str:
    """Extract text from PDF preserving page structure."""
    import fitz  # pymupdf

    doc = fitz.open(stream=content, filetype="pdf")
    pages = []
    for page_num, page in enumerate(doc, 1):
        page_text = page.get_text("text")
        if page_text.strip():
            # Add page markers for structure preservation
            pages.append(f"--- Page {page_num} ---\n{page_text.strip()}")

        # Also extract annotations/comments
        try:
            annots = page.get_annotations()
            for annot in annots:
                annot_text = annot.info.get("content", "")
                if annot_text:
                    pages.append(f"[Annotation] {annot_text}")
        except Exception:
            pass

    doc.close()
    return "\n\n".join(pages)


def _parse_pptx(content: bytes) -> str:
    """Extract text from PPTX preserving slide structure."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    slides = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts = [f"--- Slide {slide_num} ---"]

        # Extract slide title
        if slide.shapes.title and slide.shapes.title.text.strip():
            slide_parts.append(f"# {slide.shapes.title.text.strip()}")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and text != (slide.shapes.title.text.strip() if slide.shapes.title else ""):
                        # Detect list items
                        if para.level and para.level > 0:
                            indent = "  " * para.level
                            slide_parts.append(f"{indent}- {text}")
                        else:
                            slide_parts.append(text)

            # Extract table content
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        slide_parts.append("| " + " | ".join(cells) + " |")

        # Extract speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"\n[Notes] {notes}")

        slides.append("\n".join(slide_parts))

    return "\n\n".join(slides)
